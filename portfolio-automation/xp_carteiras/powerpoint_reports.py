"""Powerpoint Reports do pipeline de carteiras XP."""

from __future__ import annotations

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

from .excel_reports import _fmt_dur, _fmt_int_br, _fmt_num_br, _fmt_pct_br
from .performance import (
    _nome_col_carteira,
    _ret_ano,
    estatisticas_carteira,
    info_adicionais_lamina,
    recortar_ultimo_mes_fechado,
    tabela_performance_mensal_ano,
    tabela_retornos_acumulados,
)
from .constants import rotulo_carteira_pt


_BENCHMARK_LABELS = {
    'IBOV': 'Ibovespa',
    'IBOVESPA': 'Ibovespa',
    'SMLL': 'SMLL',
    'ISE': 'ISEE',
    'ISEE': 'ISEE',
}


def _benchmark_from_label(label):
    """Converte o rótulo do template para o nome da coluna de benchmark."""
    return _BENCHMARK_LABELS.get(str(label).strip().upper())


def _resolve_benchmark(label, available, used=None):
    """Resolve o rótulo do template para um benchmark disponível.

    Templates antigos podem trazer IBOV mesmo quando o dataframe comercial foi
    filtrado para SMLL. Nesse caso, usa o próximo benchmark disponível.
    """
    used = set() if used is None else used
    exact = _benchmark_from_label(label)
    if exact in available and exact not in used:
        return exact
    return next((benchmark for benchmark in available if benchmark not in used), None)


def _benchmark_columns(df_port):
    nome_cart = _nome_col_carteira(df_port)
    return [column for column in df_port.columns if column != nome_cart]

def _set_cell_text(cell, value):
    """Troca o texto da célula preservando a formatação do primeiro run."""
    p = cell.text_frame.paragraphs[0]
    runs = p.runs
    if runs:
        runs[0].text = str(value)
        for r in runs[1:]:                 # remove runs extras (mantém a fonte do 1º)
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text = str(value)


def _classifica_tabelas(slide):
    """Localiza cada tabela da lâmina pelo texto do cabeçalho (robusto a posição)."""
    achadas = {}
    for shp in slide.shapes:
        if not shp.has_table:
            continue
        t = shp.table
        r0c0 = t.rows[0].cells[0].text.strip()
        if r0c0 == 'Retornos anos anteriores':
            achadas['anos'] = t
        elif r0c0 == 'Retornos acumulados':
            achadas['acumulados'] = t
        elif r0c0 == 'Estatísticas':
            achadas['estatisticas'] = t
        elif r0c0 == 'Informações adicionais':
            achadas['info'] = t
        else:
            # tabela mensal "Retornos <ano>": cabeçalho começa com 'Jan'
            try:
                if t.rows[0].cells[1].text.strip() == 'Jan':
                    achadas['mensal'] = t
            except (IndexError, AttributeError):
                pass
    return achadas


def _mapa_colunas(table, nome_cart, df_port):
    """Mapeia as colunas do template para as séries presentes no dataframe."""
    header = [c.text.strip() for c in table.rows[0].cells]
    available = _benchmark_columns(df_port)
    used = set()
    col_map = {}
    for ci, h in enumerate(header):
        if ci == 0:
            continue
        if h == 'Carteira':
            col_map[ci] = nome_cart
        elif _benchmark_from_label(h) is not None:
            benchmark = _resolve_benchmark(h, available, used)
            if benchmark is None:
                continue
            col_map[ci] = benchmark
            used.add(benchmark)
            if _benchmark_from_label(h) != benchmark:
                _set_cell_text(table.rows[0].cells[ci], benchmark)
    return col_map


def _atualiza_tabela_mensal(table, df_port, portfolio, ano=None):
    t = tabela_performance_mensal_ano(df_port, ano)      # linhas cart/bench ; cols Jan..Dez(en)+ano
    ano_lbl = t.columns[-1]
    meses_en = list(t.columns[:-1])
    nome_cart = _nome_col_carteira(df_port)
    rot_cart = rotulo_carteira_pt.get(portfolio, nome_cart)
    available = _benchmark_columns(df_port)
    used = set()

    # atualiza o rótulo do ano no cabeçalho (última coluna da 1ª linha)
    cabec = table.rows[0].cells
    _set_cell_text(cabec[len(cabec) - 1], str(ano_lbl))

    for row in list(table.rows)[1:]:
        rot = row.cells[0].text.strip()
        if rot in (rot_cart, nome_cart):
            serie_col = nome_cart
        elif rot in available:
            serie_col = rot
            used.add(serie_col)
        elif _benchmark_from_label(rot) is not None:
            serie_col = _resolve_benchmark(rot, available, used)
            if serie_col is None:
                continue
            used.add(serie_col)
            if _benchmark_from_label(rot) != serie_col:
                _set_cell_text(row.cells[0], serie_col)
        else:
            continue
        vals = t.loc[serie_col]
        cells = row.cells
        for j, m in enumerate(meses_en, start=1):        # c1..c12 = meses
            _set_cell_text(cells[j], _fmt_pct_br(vals[m], 1))
        _set_cell_text(cells[len(cells) - 1], _fmt_pct_br(vals[ano_lbl], 1))


def _atualiza_tabela_anos(table, df_port):
    nome_cart = _nome_col_carteira(df_port)
    col_map = _mapa_colunas(table, nome_cart, df_port)
    for row in list(table.rows)[1:]:
        try:
            ano = int(row.cells[0].text.strip())
        except ValueError:
            continue
        for ci, serie in col_map.items():
            if serie in df_port.columns:
                _set_cell_text(row.cells[ci], _fmt_pct_br(_ret_ano(df_port[serie], ano), 1))


def _atualiza_tabela_acumulados(table, df_port):
    nome_cart = _nome_col_carteira(df_port)
    t = tabela_retornos_acumulados(df_port)              # linhas períodos ; cols cart/bench
    col_map = _mapa_colunas(table, nome_cart, df_port)
    for row in list(table.rows)[1:]:
        per = row.cells[0].text.strip()
        if per not in t.index:
            continue
        for ci, serie in col_map.items():
            if serie in t.columns:
                _set_cell_text(row.cells[ci], _fmt_pct_br(t.loc[per, serie], 1))


def _atualiza_2col_por_rotulo(table, mapa):
    """Atualiza a 2ª coluna das linhas cujo rótulo (1ª coluna) casa com o mapa."""
    for row in table.rows:
        rot = row.cells[0].text.strip()
        if rot in mapa and mapa[rot] is not None:
            _set_cell_text(row.cells[1], mapa[rot])


def _mapa_estatisticas(serie):
    e = estatisticas_carteira(serie)
    return {
        'Meses positivos':        _fmt_int_br(e['Meses positivos']),
        'Meses negativos':        _fmt_int_br(e['Meses negativos']),
        'Retorno médio mensal':   _fmt_pct_br(e['Retorno médio mensal'], 2),
        'Retorno máximo mensal':  _fmt_pct_br(e['Retorno máximo mensal'], 2),
        'Retorno mínimo mensal':  _fmt_pct_br(e['Retorno mínimo mensal'], 2),
        'Maior drawdown':         _fmt_pct_br(e['Maior drawdown'], 2),
        'Duração do maior drawdown': _fmt_dur(e['Duração do maior drawdown (meses)']),
    }


def _mapa_info(df_port, composition, serie_cdi):
    """Só os 3 campos calculados. Os fixos (corretagem, clientes) ficam intocados no PPT."""
    info = info_adicionais_lamina(df_port, composition, serie_cdi)
    return {
        'Giro médio mensal*':       _fmt_pct_br(info['Giro médio mensal'], 1),
        'Volatilidade anualizada*': _fmt_pct_br(info['Volatilidade anualizada'], 2),
        'Índice de Sharpe*':        _fmt_num_br(info['Índice de Sharpe'], 2),
    }


def _nome_serie_xml(serie_element):
    """Lê o nome original da série no XML (c:tx), para preservar a legenda."""
    tx = serie_element.find(qn('c:tx'))
    if tx is None:
        return None
    valores = [v.text for v in tx.iter(qn('c:v')) if v.text]
    return valores[0] if valores else None


def _atualiza_grafico(shape, df_port):
    """Substitui os dados do gráfico de linha, preservando estilo E nomes das séries."""
    nome_cart = _nome_col_carteira(df_port)
    chart = shape.chart

    # ordem das colunas = mesma ordem das séries do gráfico (carteira, depois benchmarks)
    cols = [nome_cart] + [c for c in df_port.columns if c != nome_cart]
    nomes_orig = [_nome_serie_xml(s._element) for s in chart.series]

    def _limpa(col):
        return [None if pd.isna(v) else float(v) for v in df_port[col].tolist()]

    cd = CategoryChartData(number_format=r'[$-416]mmm\-yy;@')
    cd.categories = [pd.Timestamp(d).to_pydatetime() for d in df_port.index]
    for i, col in enumerate(cols):
        nome_original = nomes_orig[i] if i < len(nomes_orig) else None
        if i == 0:
            nome = nome_original or col
        elif _benchmark_from_label(nome_original) == col:
            nome = nome_original
        else:
            nome = col
        cd.add_series(nome, _limpa(col))
    chart.replace_data(cd)


def atualizar_ppt(
    caminho_template,
    caminho_saida,
    df_port,
    composition,
    serie_cdi,
    portfolio,
    ano=None,
    *,
    today=None,
):
    """Preenche a lâmina somente com performance de meses já fechados."""
    df_fechado = recortar_ultimo_mes_fechado(df_port, today=today)
    prs = Presentation(caminho_template)
    for slide in prs.slides:
        tabs = _classifica_tabelas(slide)
        if 'mensal' in tabs:
            _atualiza_tabela_mensal(tabs['mensal'], df_fechado, portfolio, ano)
        if 'anos' in tabs:
            _atualiza_tabela_anos(tabs['anos'], df_fechado)
        if 'acumulados' in tabs:
            _atualiza_tabela_acumulados(tabs['acumulados'], df_fechado)
        if 'estatisticas' in tabs:
            _atualiza_2col_por_rotulo(
                tabs['estatisticas'],
                _mapa_estatisticas(df_fechado[_nome_col_carteira(df_fechado)]),
            )
        if 'info' in tabs:
            _atualiza_2col_por_rotulo(
                tabs['info'], _mapa_info(df_fechado, composition, serie_cdi)
            )
        for shp in slide.shapes:
            if shp.has_chart:
                _atualiza_grafico(shp, df_fechado)
    prs.save(caminho_saida)
    print(f"PPT atualizado: {caminho_saida}")


