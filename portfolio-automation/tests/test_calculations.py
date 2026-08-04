"""Testes unitários que não dependem dos arquivos corporativos."""

from __future__ import annotations

import unittest
from datetime import date
from unittest.mock import patch

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from xp_carteiras import email_report as email
from xp_carteiras import performance as portfolio
from xp_carteiras.accountability_reports import (
    AccountabilityPeriod,
    accountability_output_filename,
    performance_summary,
    reconcile_decomposition_total,
    resolve_accountability_period,
    waterfall_arrays,
)
from xp_carteiras.components import decomposicao_retorno, gerar_tabelas_componentes
from xp_carteiras.constants import arq_base100, arq_performance, mapa_arquivo
from xp_carteiras.email_report import EMAIL_FILENAME
from xp_carteiras.performance import (
    _df_para_lamina,
    recortar_ultimo_mes_fechado,
    tabela_performance_mensal_ano,
    tabela_retornos_acumulados,
)
from xp_carteiras.monthly_config import ACCOUNTABILITY_DISPLAY_LABELS, ACCOUNTABILITY_OUTPUT_LABELS
from xp_carteiras.powerpoint_reports import (
    _atualiza_grafico,
    _mapa_colunas,
    _nome_serie_xml,
)
from xp_carteiras.settings import Settings


class EmailCalculationsTest(unittest.TestCase):
    def setUp(self) -> None:
        self.index = pd.date_range("2025-01-01", periods=400, freq="D")
        self.series = pd.Series(np.linspace(100.0, 150.0, len(self.index)), index=self.index)

    def test_calcular_janelas_includes_all_expected_periods(self) -> None:
        result = email.calcular_janelas(self.series)

        self.assertEqual(list(result), email.WINDOW_COLUMNS)
        self.assertAlmostEqual(result["Desde Início"], 0.5)
        self.assertGreater(result["1D"], 0)

    def test_calcular_risco_returns_finite_metrics(self) -> None:
        ibov = pd.Series(np.linspace(100.0, 140.0, len(self.index)), index=self.index)
        cdi = pd.Series(np.linspace(100.0, 110.0, len(self.index)), index=self.index)

        result = email.calcular_risco(self.series, ibov, cdi)

        self.assertEqual(list(result), email.RISK_COLUMNS)
        self.assertTrue(np.isfinite(result["VOLATILIDADE"]))
        self.assertTrue(np.isfinite(result["BETA"]))
        self.assertLessEqual(result["DRAWDOWN"], 0)


class PortfolioCalculationsTest(unittest.TestCase):
    def test_daily_return_without_drift_rebalances_to_target(self) -> None:
        returns = pd.DataFrame(
            {
                "cod_ativo": ["AAA", "BBB", "AAA", "BBB"],
                "data": pd.to_datetime(["2026-01-02", "2026-01-02", "2026-01-05", "2026-01-05"]),
                "ret": [0.10, 0.00, 0.00, 0.10],
            }
        )
        composition = pd.DataFrame({"cod_ativo": ["AAA", "BBB"], "peso": [0.5, 0.5]})

        result = portfolio._ret_diario_sem_drift(returns, composition)

        np.testing.assert_allclose(result.to_numpy(), [0.05, 0.05])

    def test_daily_return_with_drift_uses_evolved_weights(self) -> None:
        returns = pd.DataFrame(
            {
                "cod_ativo": ["AAA", "BBB", "AAA", "BBB"],
                "data": pd.to_datetime(["2026-01-02", "2026-01-02", "2026-01-05", "2026-01-05"]),
                "ret": [0.10, 0.00, 0.00, 0.10],
            }
        )
        composition = pd.DataFrame({"cod_ativo": ["AAA", "BBB"], "peso": [0.5, 0.5]})

        result = portfolio._ret_diario_com_drift(returns, composition)

        self.assertAlmostEqual(result.iloc[0], 0.05)
        self.assertAlmostEqual(result.iloc[1], 0.1 * (0.5 / 1.05))


class PriorCompositionCurrentMtdTest(unittest.TestCase):
    def test_uses_prior_composition_with_current_month_returns(self) -> None:
        composition = pd.DataFrame(
            {
                "cod_ativo": ["AAA", "BBB", "CCC"],
                pd.Timestamp("2026-05-29"): [0.7, 0.3, np.nan],
                pd.Timestamp("2026-06-30"): [np.nan, 0.4, 0.6],
                pd.Timestamp("2026-07-31"): [0.2, np.nan, 0.8],
            }
        )
        prices = {
            "AAA": [90.0, 95.0, 100.0, 105.0],
            "BBB": [80.0, 90.0, 100.0, 110.0],
            "CCC": [220.0, 210.0, 200.0, 180.0],
        }
        dates = pd.to_datetime(["2026-05-29", "2026-06-30", "2026-07-31", "2026-08-14"])
        market_data = pd.DataFrame(
            [
                {"cod_ativo": ticker, "data": date, "adj_close_price": price}
                for ticker, ticker_prices in prices.items()
                for date, price in zip(dates, ticker_prices)
            ]
        )

        _, _, prior_composition_current_mtd = gerar_tabelas_componentes(
            composition,
            market_data,
            name_map={"AAA": "A", "BBB": "B", "CCC": "C"},
            sector_map={"AAA": "Setor A", "BBB": "Setor B", "CCC": "Setor C"},
        )

        result = prior_composition_current_mtd.set_index("Ticker")
        self.assertEqual(set(result.index), {"BBB", "CCC"})
        self.assertAlmostEqual(result.loc["BBB", "Peso"], 0.4)
        self.assertAlmostEqual(result.loc["CCC", "Peso"], 0.6)
        self.assertAlmostEqual(result.loc["BBB", "Desempenho no mês"], 0.10)
        self.assertAlmostEqual(result.loc["CCC", "Desempenho no mês"], -0.10)


class ReturnAttributionCompositionTest(unittest.TestCase):
    def test_month_uses_composition_in_force_before_month_end_rebalance(self) -> None:
        composition = pd.DataFrame(
            {
                "cod_ativo": ["AAA", "BBB"],
                pd.Timestamp("2026-06-30"): [0.8, 0.2],
                pd.Timestamp("2026-07-31"): [0.2, 0.8],
            }
        )
        market_data = pd.DataFrame(
            {
                "cod_ativo": ["AAA", "AAA", "BBB", "BBB"],
                "data": pd.to_datetime(["2026-06-30", "2026-07-31", "2026-06-30", "2026-07-31"]),
                "adj_close_price": [100.0, 110.0, 100.0, 100.0],
            }
        )

        result = decomposicao_retorno(
            composition,
            2026,
            7,
            market_data,
            name_map={"AAA": "A", "BBB": "B"},
            sector_map={"AAA": "Setor A", "BBB": "Setor B"},
        ).set_index("Ticker")

        self.assertAlmostEqual(result.loc["AAA", "Peso"], 0.8)
        self.assertAlmostEqual(result.loc["BBB", "Peso"], 0.2)
        self.assertAlmostEqual(result.loc["", "Contribuição"], 0.08)

class SettingsTest(unittest.TestCase):
    def test_environment_overrides_output_directory(self) -> None:
        with patch.dict("os.environ", {"XP_OUTPUT_DIR": r"C:\temp\xp-output"}):
            settings = Settings.from_env()

        self.assertEqual(str(settings.output_dir), r"C:\temp\xp-output")


class ArtifactNamesContractTest(unittest.TestCase):
    def test_external_filenames_remain_compatible(self) -> None:
        self.assertEqual(
            set(arq_base100.values()),
            {
                "top_acoes_base_100",
                "top_dividendos_base_100",
                "top_small_caps_base_100",
                "esg_base_100",
            },
        )
        self.assertEqual(
            set(arq_performance.values()),
            {
                "tab_performance_top_acoes.xlsx",
                "tab_performance_top_dividendos.xlsx",
                "tab_performance_top_small_caps.xlsx",
                "tab_performance_esg.xlsx",
            },
        )
        self.assertEqual(EMAIL_FILENAME, "email_carteiras.msg")
        self.assertEqual(
            {
                f"componentes_{name}_comp_mes_passado_mtd_atual.xlsx"
                for name in mapa_arquivo.values()
            },
            {
                "componentes_top_acoes_comp_mes_passado_mtd_atual.xlsx",
                "componentes_top_dividendos_comp_mes_passado_mtd_atual.xlsx",
                "componentes_top_small_caps_comp_mes_passado_mtd_atual.xlsx",
                "componentes_esg_comp_mes_passado_mtd_atual.xlsx",
            },
        )


class CommercialDeckBenchmarkTest(unittest.TestCase):
    def test_current_partial_month_is_excluded_from_year_and_last_12_months(self) -> None:
        portfolio_name = "Carteira - TOP Ações XP"
        frame = pd.DataFrame(
            {
                portfolio_name: [80.0, 100.0, 100.0, 110.0, 220.0],
                "Ibovespa": [80.0, 100.0, 100.0, 108.0, 216.0],
            },
            index=pd.to_datetime(
                ["2025-07-31", "2025-12-31", "2026-06-30", "2026-07-31", "2026-08-04"]
            ),
        )

        closed = recortar_ultimo_mes_fechado(frame, today=date(2026, 8, 4))
        annual = tabela_performance_mensal_ano(closed, 2026)
        accumulated = tabela_retornos_acumulados(closed)

        self.assertEqual(closed.index.max(), pd.Timestamp("2026-07-31"))
        self.assertAlmostEqual(annual.loc[portfolio_name, "2026"], 0.10)
        self.assertAlmostEqual(
            accumulated.loc["Últimos 12 meses", portfolio_name], 110 / 80 - 1
        )

    def test_small_caps_uses_smll_instead_of_ibovespa(self) -> None:
        portfolio_name = "Carteira - TOP SMALL CAPS XP"
        index = pd.to_datetime(["2026-01-30", "2026-02-27"])
        result_frames = {
            portfolio_name: pd.DataFrame(
                {
                    portfolio_name: [100.0, 102.0],
                    "Ibovespa": [100.0, 101.0],
                    "SMLL": [100.0, 103.0],
                },
                index=index,
            )
        }

        commercial = _df_para_lamina(portfolio_name, result_frames)

        self.assertEqual(list(commercial.columns), [portfolio_name, "SMLL"])

    def test_chart_relabels_legacy_ibovespa_series_as_smll(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        chart_data = CategoryChartData()
        chart_data.categories = ["jan-26", "fev-26"]
        chart_data.add_series("Carteira", [100.0, 102.0])
        chart_data.add_series("Ibovespa", [100.0, 101.0])
        shape = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            Inches(1), Inches(1), Inches(8), Inches(4),
            chart_data,
        )
        index = pd.to_datetime(["2026-01-30", "2026-02-27"])
        commercial = pd.DataFrame(
            {
                "Carteira - TOP SMALL CAPS XP": [100.0, 102.0],
                "SMLL": [100.0, 103.0],
            },
            index=index,
        )

        _atualiza_grafico(shape, commercial)

        series_names = [_nome_serie_xml(series._element) for series in shape.chart.series]
        self.assertEqual(series_names, ["Carteira", "SMLL"])

    def test_table_relabels_legacy_ibov_header_as_smll(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        table = slide.shapes.add_table(
            2, 3,
            Inches(1), Inches(1), Inches(8), Inches(2),
        ).table
        table.rows[0].cells[0].text = "Período"
        table.rows[0].cells[1].text = "Carteira"
        table.rows[0].cells[2].text = "IBOV"
        portfolio_name = "Carteira - TOP SMALL CAPS XP"
        commercial = pd.DataFrame(
            {
                portfolio_name: [100.0, 102.0],
                "SMLL": [100.0, 103.0],
            },
            index=pd.to_datetime(["2026-01-30", "2026-02-27"]),
        )

        column_map = _mapa_colunas(table, portfolio_name, commercial)

        self.assertEqual(column_map, {1: portfolio_name, 2: "SMLL"})
        self.assertEqual(table.rows[0].cells[2].text, "SMLL")


class AccountabilityReportTest(unittest.TestCase):
    def test_current_partial_month_uses_previous_closed_month(self) -> None:
        frame = pd.DataFrame(
            {"Carteira - TOP Ações XP": [100.0, 101.0]},
            index=pd.to_datetime(["2026-07-31", "2026-08-04"]),
        )

        period = resolve_accountability_period(frame, today=date(2026, 8, 4))

        self.assertEqual(period, AccountabilityPeriod(2026, 7, 2026, 8))

    def test_lagged_data_uses_latest_available_month(self) -> None:
        frame = pd.DataFrame(
            {"Carteira - TOP Ações XP": [100.0, 102.0]},
            index=pd.to_datetime(["2026-06-30", "2026-07-31"]),
        )

        period = resolve_accountability_period(frame, today=date(2026, 8, 4))

        self.assertEqual(period, AccountabilityPeriod(2026, 7, 2026, 8))

    def test_performance_summary_closes_at_reference_month(self) -> None:
        frame = pd.DataFrame(
            {
                "Carteira - TOP Ações XP": [100.0, 110.0, 132.0, 150.0],
                "Ibovespa": [100.0, 105.0, 115.5, 120.0],
            },
            index=pd.to_datetime(["2025-06-30", "2025-12-31", "2026-05-31", "2026-06-30"]),
        )

        summary = performance_summary(frame, 2026, 6)

        self.assertAlmostEqual(summary.loc["Carteira - TOP Ações XP", "Mês"], 150 / 132 - 1)
        self.assertAlmostEqual(summary.loc["Carteira - TOP Ações XP", "Acumulado"], 150 / 110 - 1)
        self.assertAlmostEqual(summary.loc["Carteira - TOP Ações XP", "12 meses"], 150 / 100 - 1)

    def test_performance_summary_ignores_current_month_even_if_present(self) -> None:
        frame = pd.DataFrame(
            {"Carteira - TOP Ações XP": [100.0, 110.0, 220.0]},
            index=pd.to_datetime(["2026-06-30", "2026-07-31", "2026-08-04"]),
        )

        summary = performance_summary(frame, 2026, 7)

        self.assertAlmostEqual(summary.loc["Carteira - TOP Ações XP", "Mês"], 0.10)
        self.assertAlmostEqual(summary.loc["Carteira - TOP Ações XP", "Acumulado"], 0.10)

    def test_waterfall_keeps_total_and_all_contributions(self) -> None:
        categories, _, high, low, total, ordered = waterfall_arrays(
            ["AAA", "BBB", "CCC"], [0.02, -0.01, 0.005], 0.015
        )

        self.assertEqual(categories[-1], "Carteira")
        self.assertAlmostEqual(sum(ordered), 0.015)
        self.assertEqual(sum(value is not None for value in high), 2)
        self.assertEqual(sum(value is not None for value in low), 1)
        self.assertAlmostEqual(total[-1], 0.015)

    def test_output_filename_contains_report_month(self) -> None:
        period = AccountabilityPeriod(2026, 7, 2026, 8)

        filename = accountability_output_filename("Top Ações", period)

        self.assertEqual(filename, "Prestação de Contas - Top Ações - Agosto 2026.pptx")

    def test_three_portfolios_are_enabled(self) -> None:
        expected = {
            "Carteira - TOP Ações XP",
            "Carteira - TOP DIVIDENDOS XP",
            "Carteira - TOP SMALL CAPS XP",
        }

        self.assertEqual(set(ACCOUNTABILITY_OUTPUT_LABELS), expected)
        self.assertEqual(set(ACCOUNTABILITY_DISPLAY_LABELS), expected)

    def test_small_reconciliation_difference_becomes_explicit_bar(self) -> None:
        decomposition = pd.DataFrame(
            [
                {"Companhia": "A", "Ticker": "AAA", "Setor": "S", "Peso": 1.0,
                 "Retorno no mês": 0.01, "Contribuição": 0.01},
                {"Companhia": "Carteira (total)", "Ticker": "", "Setor": "", "Peso": 1.0,
                 "Retorno no mês": np.nan, "Contribuição": 0.01},
            ]
        )

        result = reconcile_decomposition_total(decomposition, 0.0102)

        self.assertIn("Ajuste", set(result["Ticker"]))
        self.assertAlmostEqual(result.loc[result["Ticker"] == "", "Contribuição"].iloc[0], 0.0102)

    def test_large_reconciliation_difference_stops_generation(self) -> None:
        decomposition = pd.DataFrame(
            [
                {"Companhia": "A", "Ticker": "AAA", "Contribuição": 0.046},
                {"Companhia": "Carteira (total)", "Ticker": "", "Contribuição": 0.046},
            ]
        )

        with self.assertRaisesRegex(ValueError, "não confere"):
            reconcile_decomposition_total(decomposition, 0.008)


if __name__ == "__main__":
    unittest.main()
